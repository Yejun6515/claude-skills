#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
doc-verify / extract.py

Excel(.xlsx/.xls) · PDF · Word(.docx) · PowerPoint(.pptx) 를 하나의 공통 구조로 뽑는다.
verify.py 가 이 구조 위에서 검증을 돌린다. 단독 실행하면 추출 결과를 눈으로 확인할 수 있다.

공통 구조:
  {
    "path": str, "type": "xlsx|xls|pdf|docx|pptx",
    "tables": [ {"loc": "Sheet1" | "p.3" | "slide 2", "rows": [[cell,...], ...],
                 "formulas": {"r,c": "=SUM(...)"} } ],
    "texts":  [ {"loc": ..., "text": ...} ],
    "revisions": {"tracked_insert": n, "tracked_delete": n, "comments": n},
  }

셀 값은 전부 문자열로 담는다(표시값 기준). 숫자 파싱은 verify.py 의 parse_num 이 한다.
"""
import sys, os, re, json, zipfile, warnings

warnings.filterwarnings("ignore")   # openpyxl DrawingML/DataValidation 경고가 stderr를 덮는다

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8")
    sys.stderr.reconfigure(encoding="utf-8")


def _s(v):
    """셀 값을 표시 문자열로. None -> '' """
    if v is None:
        return ""
    if isinstance(v, float) and v == int(v):
        return str(int(v))
    return str(v).strip()


# --------------------------------------------------------------------------- xlsx
def read_xlsx(path):
    from openpyxl import load_workbook
    wb_v = load_workbook(path, data_only=True)   # 캐시된 계산값
    wb_f = load_workbook(path, data_only=False)  # 수식 원문
    tables, texts = [], []
    comments = 0
    for ws_v in wb_v.worksheets:
        if ws_v.sheet_state != "visible":
            continue
        ws_f = wb_f[ws_v.title]
        rows, formulas = [], {}
        for r in range(1, (ws_v.max_row or 0) + 1):
            row = []
            for c in range(1, (ws_v.max_column or 0) + 1):
                cv = ws_v.cell(row=r, column=c)
                cf = ws_f.cell(row=r, column=c)
                row.append(_s(cv.value))
                if isinstance(cf.value, str) and cf.value.startswith("="):
                    formulas["%d,%d" % (r, c)] = cf.value
                if cf.comment is not None:
                    comments += 1
            rows.append(row)
        tables.append({"loc": ws_v.title, "rows": rows, "formulas": formulas})
        for r, row in enumerate(rows, 1):
            line = " | ".join(x for x in row if x)
            if line:
                texts.append({"loc": "%s!r%d" % (ws_v.title, r), "text": line})
    return tables, texts, {"tracked_insert": 0, "tracked_delete": 0, "comments": comments}


# --------------------------------------------------------------------------- xls
def read_xls(path):
    import xlrd
    book = xlrd.open_workbook(path)
    tables, texts = [], []
    for sh in book.sheets():
        rows = []
        for r in range(sh.nrows):
            rows.append([_s(sh.cell_value(r, c)) for c in range(sh.ncols)])
        tables.append({"loc": sh.name, "rows": rows, "formulas": {}})
        for r, row in enumerate(rows, 1):
            line = " | ".join(x for x in row if x)
            if line:
                texts.append({"loc": "%s!r%d" % (sh.name, r), "text": line})
    return tables, texts, {"tracked_insert": 0, "tracked_delete": 0, "comments": 0}


# --------------------------------------------------------------------------- pdf
def read_pdf(path):
    import pdfplumber
    tables, texts = [], []
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            txt = page.extract_text() or ""
            for ln, line in enumerate(txt.splitlines(), 1):
                if line.strip():
                    texts.append({"loc": "p.%d L%d" % (i, ln), "text": line.strip()})
            for t, tbl in enumerate(page.extract_tables() or [], 1):
                rows = [[_s(c) for c in row] for row in tbl]
                tables.append({"loc": "p.%d table%d" % (i, t), "rows": rows, "formulas": {}})
    return tables, texts, {"tracked_insert": 0, "tracked_delete": 0, "comments": 0}


# --------------------------------------------------------------------------- docx
def read_docx(path):
    import docx
    d = docx.Document(path)
    texts, tables = [], []
    for i, p in enumerate(d.paragraphs, 1):
        if p.text.strip():
            texts.append({"loc": "para %d" % i, "text": p.text.strip()})
    for ti, t in enumerate(d.tables, 1):
        rows = [[_s(c.text) for c in row.cells] for row in t.rows]
        tables.append({"loc": "table %d" % ti, "rows": rows, "formulas": {}})
        for r, row in enumerate(rows, 1):
            line = " | ".join(x for x in row if x)
            if line:
                texts.append({"loc": "table %d r%d" % (ti, r), "text": line})
    rev = _ooxml_revisions(path, ["word/document.xml"], "word/comments.xml")
    return tables, texts, rev


# --------------------------------------------------------------------------- pptx
def read_pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    texts, tables = [], []
    for i, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    line = "".join(r.text for r in p.runs).strip()
                    if line:
                        texts.append({"loc": "slide %d" % i, "text": line})
            if getattr(shape, "has_table", False) and shape.has_table:
                rows = [[_s(c.text) for c in row.cells] for row in shape.table.rows]
                tables.append({"loc": "slide %d table" % i, "rows": rows, "formulas": {}})
                for r, row in enumerate(rows, 1):
                    line = " | ".join(x for x in row if x)
                    if line:
                        texts.append({"loc": "slide %d table r%d" % (i, r), "text": line})
    rev = _ooxml_revisions(path, [], None, comment_glob="comment")
    return tables, texts, rev


# --------------------------------------------------------------------------- OOXML 개정 흔적
def _ooxml_revisions(path, doc_parts, comments_part, comment_glob=None):
    """track changes(w:ins/w:del)·코멘트 잔존을 zip 레벨에서 센다."""
    ins = dele = com = 0
    try:
        with zipfile.ZipFile(path) as z:
            names = z.namelist()
            for part in doc_parts:
                if part in names:
                    xml = z.read(part).decode("utf-8", "ignore")
                    ins += len(re.findall(r"<w:ins[ >]", xml))
                    dele += len(re.findall(r"<w:del[ >]", xml))
                    com += len(re.findall(r"commentReference", xml))
            if comments_part and comments_part in names:
                xml = z.read(comments_part).decode("utf-8", "ignore")
                com = max(com, len(re.findall(r"<w:comment[ >]", xml)))
            if comment_glob:
                for n in names:
                    if comment_glob in n.lower() and n.endswith(".xml"):
                        xml = z.read(n).decode("utf-8", "ignore")
                        com += len(re.findall(r"<p:cm[ >]|<pc:cm[ >]", xml))
    except Exception:
        pass
    return {"tracked_insert": ins, "tracked_delete": dele, "comments": com}


# --------------------------------------------------------------------------- entry
READERS = {
    ".xlsx": read_xlsx, ".xlsm": read_xlsx, ".xls": read_xls,
    ".pdf": read_pdf, ".docx": read_docx, ".pptx": read_pptx,
}


def extract(path):
    ext = os.path.splitext(path)[1].lower()
    if ext not in READERS:
        raise SystemExit("지원하지 않는 형식: %s (지원: %s)" % (ext, ", ".join(sorted(READERS))))
    tables, texts, rev = READERS[ext](path)
    return {"path": path, "type": ext.lstrip("."), "tables": tables,
            "texts": texts, "revisions": rev}


def main():
    if len(sys.argv) < 2:
        raise SystemExit("usage: extract.py <file> [--json]")
    doc = extract(sys.argv[1])
    if "--json" in sys.argv:
        print(json.dumps(doc, ensure_ascii=False, indent=1))
        return
    print("파일 : %s (%s)" % (doc["path"], doc["type"]))
    print("표 %d개 · 텍스트 %d줄 · 개정흔적 %s" % (len(doc["tables"]), len(doc["texts"]), doc["revisions"]))
    for t in doc["tables"]:
        print("\n=== TABLE [%s] %d행 ===" % (t["loc"], len(t["rows"])))
        for r, row in enumerate(t["rows"][:60], 1):
            if any(x for x in row):
                print("%3d | %s" % (r, " | ".join(row)))
        if len(t["rows"]) > 60:
            print("... (%d행 생략)" % (len(t["rows"]) - 60))


if __name__ == "__main__":
    main()
